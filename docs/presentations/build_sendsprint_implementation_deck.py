"""Build bilingual SendSprint implementation decks as PPTX, PDF, and PNG previews."""
# ruff: noqa: E501

from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path
from textwrap import wrap

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR
from pptx.util import Inches, Pt
from reportlab.lib.pagesizes import landscape
from reportlab.lib.units import inch
from reportlab.pdfgen import canvas

ROOT = Path(__file__).resolve().parents[2]
OUT = ROOT / "docs" / "presentations"
PREVIEWS = OUT / "preview"
ASSETS = ROOT / "docs" / "assets"
VIDEO = ROOT / "video" / "preview"

WIDE = (13.333333, 7.5)
PX = (1920, 1080)

COLORS = {
    "bg": "#06111f",
    "bg2": "#0b1b2d",
    "text": "#eef6ff",
    "muted": "#a9b7c7",
    "cyan": "#22d3ee",
    "green": "#34d399",
    "gold": "#fbbf24",
    "red": "#fb7185",
    "line": "#24445e",
    "panel": "#10243a",
}


@dataclass(frozen=True)
class Slide:
    title: str
    kicker: str
    body: str
    bullets: tuple[str, ...] = ()
    image: Path | None = None
    image_label: str | None = None
    accent: str = "cyan"


COPY: dict[str, dict[str, object]] = {
    "en": {
        "filename": "sendsprint-implementation-en",
        "deck_title": "SendSprint Implementation",
        "subtitle": "From scattered sprint work to validated PR delivery",
        "date": "May 2026",
        "slides": (
            Slide(
                "SendSprint\nImplementation",
                "PRODUCTIVITY MULTIPLIER",
                "A practical execution layer that reads Jira/Azure DevOps, plans the delivery, validates safety, implements, tests, and opens PRs.",
                ("Preflight", "Dry-run", "Resume state", "Evidence-driven PRs"),
                ASSETS / "sendsprint-productivity-engine.png",
                "SendSprint delivery engine",
            ),
            Slide(
                "Pain Before SendSprint",
                "WHY THIS MATTERS",
                "Sprint delivery often fails in the spaces between tools, not because engineers lack skill.",
                (
                    "Tasks sit in Jira/Azure while context is rediscovered manually.",
                    "Branches, tests, evidence, and PR text are repeated by hand.",
                    "Managers ask for status because the delivery flow is invisible.",
                ),
                ASSETS / "sendsprint-productivity-before-after.png",
                "Before/after productivity contrast",
                "red",
            ),
            Slide(
                "What Changed",
                "IMPLEMENTATION SCOPE",
                "We turned the workflow into a controlled, recoverable, auditable loop.",
                (
                    "Azure DevOps MCP installer and Jira/Azure core guide.",
                    "Default branch convention: feature/{number}-{title}.",
                    "Work-item hierarchy safety: invalid Issue -> Task becomes Related.",
                    "Generated task descriptions no longer mention internal tooling.",
                ),
                VIDEO / "sendsprint-before-after-poster-en.png",
                "Before/after explainer poster",
                "gold",
            ),
            Slide(
                "The New Delivery Flow",
                "END-TO-END LOOP",
                "SendSprint converts sprint items into isolated delivery lanes that end in PRs.",
                (
                    "Read sprint through MCP -> API -> browser fallback.",
                    "Map architecture and generate missing baseline docs.",
                    "Create branch/worktree, install, build, lint, test, and scan.",
                    "Commit, push, create PR, validate PR metadata, and review diff.",
                ),
                None,
                None,
                "cyan",
            ),
            Slide(
                "Preflight + Dry-Run",
                "SAFETY BEFORE MUTATION",
                "The team can see what will happen before branches, work items, files, or PRs change.",
                (
                    "Preflight validates transport, credentials, repo health, sprint read, and link safety.",
                    "Dry-run shows item, repo, branch, PR target, confidence, and routing reason.",
                    "Low-confidence routes are visible before execution.",
                ),
                None,
                None,
                "green",
            ),
            Slide(
                "Resume State + Idempotency",
                "NO DUPLICATE DELIVERY",
                "Interrupted runs can resume without recreating the same delivery work.",
                (
                    "Run state persists in .sendsprint/runs/<run-id>.json.",
                    "Completed item/repo pairs are skipped on retry.",
                    "Failed pairs keep the failure reason for investigation.",
                ),
                None,
                None,
                "gold",
            ),
            Slide(
                "Governance Built In",
                "AZURE/JIRA SAFETY",
                "The implementation captures hard-learned operational rules as code and documentation.",
                (
                    "No automatic assignee unless inherited or explicitly requested.",
                    "Parent-child links are validated before and after planning.",
                    "Jira/Azure stable rules live in .specs/integrations.",
                    "MCP remains the source of truth for live tenant state.",
                ),
                None,
                None,
                "cyan",
            ),
            Slide(
                "Quality Gates",
                "VALIDATED DELIVERY",
                "Every implementation change was tested before shipping.",
                (
                    "Python: ruff, mypy, and 175 pytest tests passing.",
                    "Remotion: TypeScript typecheck and rendered MP4 outputs.",
                    "README: images, posters, and video links included.",
                    "Version bumped to 0.11.0 with changelog entry.",
                ),
                VIDEO / "sendsprint-before-after-poster-en.png",
                "Rendered video poster",
                "green",
            ),
            Slide(
                "How a Company Uses It",
                "OPERATING MODEL",
                "Adoption becomes a simple, repeatable command sequence.",
                (
                    'sendsprint preflight azuredevops "Team\\\\Sprint 12" --workspace workspace.yaml',
                    'sendsprint run azuredevops "Team\\\\Sprint 12" --dry-run',
                    'sendsprint run azuredevops "Team\\\\Sprint 12" --run-id sprint-12',
                    "Open PRs with evidence and review metadata.",
                ),
                None,
                None,
                "gold",
            ),
            Slide(
                "Outcome",
                "WHAT PRODUCTIVITY FEELS LIKE",
                "Less coordination debt, fewer backlog mistakes, safer automation, and a repeatable path from sprint to develop.",
                (
                    "Engineers spend more time solving problems.",
                    "Managers get visible delivery state.",
                    "PRs arrive with evidence by default.",
                    "The same process scales across teams and repos.",
                ),
                ASSETS / "sendsprint-productivity-engine.png",
                "Productivity engine",
                "green",
            ),
        ),
    },
    "pt-BR": {
        "filename": "sendsprint-implementation-pt-BR",
        "deck_title": "Implementacao do SendSprint",
        "subtitle": "Da sprint espalhada para PRs validados",
        "date": "Maio de 2026",
        "slides": (
            Slide(
                "Implementacao\ndo SendSprint",
                "MULTIPLICADOR DE PRODUTIVIDADE",
                "Uma camada de execucao que le Jira/Azure DevOps, planeja a entrega, valida seguranca, implementa, testa e abre PRs.",
                ("Preflight", "Dry-run", "Estado resumivel", "PRs com evidencia"),
                ASSETS / "sendsprint-productivity-engine.png",
                "Motor de entrega SendSprint",
            ),
            Slide(
                "Dor Antes do SendSprint",
                "POR QUE ISSO IMPORTA",
                "A entrega da sprint normalmente falha nos espacos entre ferramentas, nao por falta de competencia do time.",
                (
                    "Tasks ficam no Jira/Azure enquanto o contexto e redescoberto manualmente.",
                    "Branches, testes, evidencias e texto de PR sao repetidos a mao.",
                    "Gestores perguntam status porque o fluxo de entrega nao esta visivel.",
                ),
                ASSETS / "sendsprint-productivity-before-after.png",
                "Contraste antes/depois",
                "red",
            ),
            Slide(
                "O Que Mudou",
                "ESCOPO IMPLEMENTADO",
                "Transformamos o workflow em um loop controlado, recuperavel e auditavel.",
                (
                    "Instalador MCP do Azure DevOps e guia core Jira/Azure.",
                    "Convencao padrao de branch: feature/{number}-{title}.",
                    "Seguranca de hierarquia: Issue -> Task invalido vira Related.",
                    "Descricoes geradas nao mencionam ferramenta interna.",
                ),
                VIDEO / "sendsprint-before-after-poster-pt.png",
                "Poster do video antes/depois",
                "gold",
            ),
            Slide(
                "O Novo Fluxo de Entrega",
                "LOOP PONTA-A-PONTA",
                "SendSprint transforma itens da sprint em esteiras isoladas que terminam em PR.",
                (
                    "Le sprint por MCP -> API -> fallback browser.",
                    "Mapeia arquitetura e gera baseline quando falta.",
                    "Cria branch/worktree, instala, builda, roda lint, testes e scan.",
                    "Comita, faz push, cria PR, valida metadados e revisa diff.",
                ),
                None,
                None,
                "cyan",
            ),
            Slide(
                "Preflight + Dry-Run",
                "SEGURANCA ANTES DE ALTERAR",
                "O time ve o que vai acontecer antes de alterar branches, work items, arquivos ou PRs.",
                (
                    "Preflight valida transporte, credenciais, repo, sprint e links.",
                    "Dry-run mostra item, repo, branch, target, confianca e motivo.",
                    "Rotas com baixa confianca ficam visiveis antes da execucao.",
                ),
                None,
                None,
                "green",
            ),
            Slide(
                "Estado Resumivel + Idempotencia",
                "SEM ENTREGA DUPLICADA",
                "Execucoes interrompidas podem continuar sem recriar o mesmo trabalho.",
                (
                    "Estado fica em .sendsprint/runs/<run-id>.json.",
                    "Pares item/repo concluidos sao ignorados no retry.",
                    "Falhas preservam o motivo para investigacao.",
                ),
                None,
                None,
                "gold",
            ),
            Slide(
                "Governanca Embutida",
                "SEGURANCA AZURE/JIRA",
                "A implementacao transforma regras operacionais aprendidas em codigo e documentacao.",
                (
                    "Sem assignee automatico, exceto herdado ou pedido explicitamente.",
                    "Links parent-child sao validados antes e depois do planejamento.",
                    "Regras estaveis de Jira/Azure vivem em .specs/integrations.",
                    "MCP continua sendo fonte de verdade para estado vivo do tenant.",
                ),
                None,
                None,
                "cyan",
            ),
            Slide(
                "Gates de Qualidade",
                "ENTREGA VALIDADA",
                "Mudancas sao testadas antes do push.",
                (
                    "Python: ruff, mypy e 175 testes pytest passando.",
                    "Remotion: typecheck TypeScript e MP4s renderizados.",
                    "README: imagens, posters e links de video incluidos.",
                    "Versao 0.11.0 com entrada no changelog.",
                ),
                VIDEO / "sendsprint-before-after-poster-pt.png",
                "Poster renderizado",
                "green",
            ),
            Slide(
                "Como a Empresa Usa",
                "MODELO OPERACIONAL",
                "A adocao vira uma sequencia simples e repetivel de comandos.",
                (
                    'sendsprint preflight azuredevops "Team\\\\Sprint 12" --workspace workspace.yaml',
                    'sendsprint run azuredevops "Team\\\\Sprint 12" --dry-run',
                    'sendsprint run azuredevops "Team\\\\Sprint 12" --run-id sprint-12',
                    "PRs abertos com evidencia e metadados de revisao.",
                ),
                None,
                None,
                "gold",
            ),
            Slide(
                "Resultado",
                "COMO A PRODUTIVIDADE APARECE",
                "Menos divida de coordenacao, menos erros de backlog, automacao mais segura e caminho repetivel da sprint para develop.",
                (
                    "Engenheiros passam mais tempo resolvendo problemas.",
                    "Gestores enxergam o estado real da entrega.",
                    "PRs chegam com evidencia por padrao.",
                    "O mesmo processo escala por times e repos.",
                ),
                ASSETS / "sendsprint-productivity-engine.png",
                "Motor de produtividade",
                "green",
            ),
        ),
    },
}


def hex_rgb(value: str) -> RGBColor:
    value = value.lstrip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def rgb_tuple(value: str) -> tuple[int, int, int]:
    value = value.lstrip("#")
    return int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16)


def add_textbox(slide, x, y, w, h, text, size=24, color="text", bold=False, font="Aptos"):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    box.text_frame.margin_left = Inches(0)
    box.text_frame.margin_right = Inches(0)
    box.text_frame.margin_top = Inches(0)
    box.text_frame.margin_bottom = Inches(0)
    box.text_frame.vertical_anchor = MSO_ANCHOR.TOP
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.font.name = font
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = hex_rgb(COLORS[color])
    return box


def add_picture_contain(slide, path: Path, x: float, y: float, w: float, h: float):
    with Image.open(path) as img:
        iw, ih = img.size
    scale = min(w / iw, h / ih)
    pw = iw * scale
    ph = ih * scale
    px = x + (w - pw) / 2
    py = y + (h - ph) / 2
    return slide.shapes.add_picture(str(path), Inches(px), Inches(py), Inches(pw), Inches(ph))


def build_pptx(lang: str, cfg: dict[str, object]) -> Path:
    prs = Presentation()
    prs.slide_width = Inches(WIDE[0])
    prs.slide_height = Inches(WIDE[1])
    blank = prs.slide_layouts[6]
    slides: tuple[Slide, ...] = cfg["slides"]  # type: ignore[assignment]

    for idx, item in enumerate(slides, start=1):
        slide = prs.slides.add_slide(blank)
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = hex_rgb(COLORS["bg"])

        accent = COLORS[item.accent]
        shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
        shape.fill.solid()
        shape.fill.fore_color.rgb = hex_rgb(COLORS["bg"])
        shape.line.fill.background()

        slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(0.13)).fill.solid()
        slide.shapes[-1].fill.fore_color.rgb = hex_rgb(accent)
        slide.shapes[-1].line.fill.background()

        add_textbox(slide, 0.65, 0.46, 4.8, 0.26, item.kicker, 10, item.accent, True, "Aptos")
        title_h = 1.88 if idx == 1 else 1.05
        body_y = 2.78 if idx == 1 else 1.83
        add_textbox(
            slide, 0.65, 0.82, 7.0, title_h, item.title, 34 if idx > 1 else 39, "text", True
        )
        add_textbox(slide, 0.68, body_y, 6.45, 0.78, item.body, 17, "muted")
        add_textbox(slide, 12.15, 0.48, 0.5, 0.2, f"{idx:02d}", 11, "muted", True)

        if item.image:
            slide.shapes.add_shape(
                1, Inches(7.55), Inches(0.9), Inches(5.1), Inches(4.0)
            ).fill.solid()
            slide.shapes[-1].fill.fore_color.rgb = hex_rgb(COLORS["panel"])
            slide.shapes[-1].line.color.rgb = hex_rgb(COLORS["line"])
            add_picture_contain(slide, item.image, 7.75, 1.05, 4.7, 3.7)
            if item.image_label:
                add_textbox(slide, 7.75, 5.02, 4.6, 0.26, item.image_label, 11, "muted", False)
            bullet_x, bullet_y, bullet_w = 0.72, 4.12 if idx == 1 else 3.0, 6.3
        else:
            bullet_x, bullet_y, bullet_w = 1.0, 4.12 if idx == 1 else 3.0, 10.8

        for b_idx, bullet in enumerate(item.bullets):
            y = bullet_y + b_idx * 0.68
            dot = slide.shapes.add_shape(
                9, Inches(bullet_x), Inches(y + 0.08), Inches(0.13), Inches(0.13)
            )
            dot.fill.solid()
            dot.fill.fore_color.rgb = hex_rgb(accent)
            dot.line.fill.background()
            add_textbox(slide, bullet_x + 0.28, y, bullet_w, 0.44, bullet, 16, "text")

        footer = "SendSprint implementation | sprint-to-PR execution layer"
        if lang == "pt-BR":
            footer = "Implementacao SendSprint | camada de execucao sprint-para-PR"
        add_textbox(slide, 0.65, 7.05, 6.0, 0.18, footer, 9, "muted")

    path = OUT / f"{cfg['filename']}.pptx"
    prs.save(path)
    return path


def draw_wrapped(
    draw: ImageDraw.ImageDraw, xy, text: str, font, fill, width_chars: int, line_gap: int = 8
):
    x, y = xy
    for line in wrap(text, width_chars):
        draw.text((x, y), line, font=font, fill=fill)
        y += font.size + line_gap
    return y


def make_preview_image(item: Slide, idx: int, lang: str, out: Path) -> None:
    img = Image.new("RGB", PX, rgb_tuple(COLORS["bg"]))
    draw = ImageDraw.Draw(img)
    font_dir = Path("C:/Windows/Fonts")
    title_font = ImageFont.truetype(str(font_dir / "segoeuib.ttf"), 82 if idx > 1 else 96)
    body_font = ImageFont.truetype(str(font_dir / "segoeui.ttf"), 38)
    bullet_font = ImageFont.truetype(str(font_dir / "segoeui.ttf"), 34)
    small_font = ImageFont.truetype(str(font_dir / "consola.ttf"), 24)
    accent = rgb_tuple(COLORS[item.accent])

    draw.rectangle((0, 0, PX[0], 22), fill=accent)
    draw.text((92, 72), item.kicker, font=small_font, fill=accent)
    title_y = 135
    for line in item.title.splitlines():
        draw.text((92, title_y), line, font=title_font, fill=rgb_tuple(COLORS["text"]))
        title_y += title_font.size + 12
    draw_wrapped(
        draw, (98, 455 if idx == 1 else 290), item.body, body_font, rgb_tuple(COLORS["muted"]), 58
    )
    draw.text((1765, 72), f"{idx:02d}", font=small_font, fill=rgb_tuple(COLORS["muted"]))

    bullet_x, bullet_y, bullet_width = 110, 630 if idx == 1 else 460, 58
    if item.image:
        panel = (1088, 150, 1815, 720)
        draw.rounded_rectangle(
            panel,
            radius=44,
            fill=rgb_tuple(COLORS["panel"]),
            outline=rgb_tuple(COLORS["line"]),
            width=2,
        )
        with Image.open(item.image).convert("RGB") as asset:
            asset.thumbnail((680, 500))
            ax = panel[0] + (panel[2] - panel[0] - asset.width) // 2
            ay = panel[1] + (panel[3] - panel[1] - asset.height) // 2
            img.paste(asset, (ax, ay))
        if item.image_label:
            draw.text(
                (1115, 750), item.image_label, font=small_font, fill=rgb_tuple(COLORS["muted"])
            )
        bullet_width = 52

    for bullet in item.bullets:
        draw.ellipse((bullet_x, bullet_y + 12, bullet_x + 18, bullet_y + 30), fill=accent)
        bullet_y = draw_wrapped(
            draw,
            (bullet_x + 42, bullet_y),
            bullet,
            bullet_font,
            rgb_tuple(COLORS["text"]),
            bullet_width,
            5,
        )
        bullet_y += 26

    footer = "SendSprint implementation | sprint-to-PR execution layer"
    if lang == "pt-BR":
        footer = "Implementacao SendSprint | camada de execucao sprint-para-PR"
    draw.text((92, 1010), footer, font=small_font, fill=rgb_tuple(COLORS["muted"]))
    img.save(out)


def build_pdf_and_previews(lang: str, cfg: dict[str, object]) -> tuple[Path, list[Path]]:
    slides: tuple[Slide, ...] = cfg["slides"]  # type: ignore[assignment]
    deck_dir = PREVIEWS / str(cfg["filename"])
    deck_dir.mkdir(parents=True, exist_ok=True)
    preview_paths: list[Path] = []

    for idx, item in enumerate(slides, start=1):
        preview = deck_dir / f"slide-{idx:02d}.png"
        make_preview_image(item, idx, lang, preview)
        preview_paths.append(preview)

    pdf_path = OUT / f"{cfg['filename']}.pdf"
    c = canvas.Canvas(str(pdf_path), pagesize=landscape((13.333 * inch, 7.5 * inch)))
    for preview in preview_paths:
        c.drawImage(str(preview), 0, 0, width=13.333 * inch, height=7.5 * inch)
        c.showPage()
    c.save()
    return pdf_path, preview_paths


def build_all() -> None:
    OUT.mkdir(parents=True, exist_ok=True)
    PREVIEWS.mkdir(parents=True, exist_ok=True)
    for lang, cfg in COPY.items():
        pptx = build_pptx(lang, cfg)
        pdf, previews = build_pdf_and_previews(lang, cfg)
        print(f"{lang}: {pptx}")
        print(f"{lang}: {pdf}")
        print(f"{lang}: {previews[0]} ... {len(previews)} previews")


if __name__ == "__main__":
    build_all()
